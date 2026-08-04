#!/usr/bin/env python3
"""
nust_crawler.py
================
Full-site crawler + scraper for https://nust.edu.pk (and its subdomains),
built for feeding a RAG (Retrieval-Augmented Generation) pipeline.

Design goals
------------
1. NO LLM anywhere in the pipeline (as requested) — pure heuristic /
   rule-based extraction:
     - HTML  -> crawl4ai's LXML scraping strategy + PruningContentFilter
                (a non-LLM, heuristic "boilerplate remover") -> clean Markdown.
     - Tables -> crawl4ai's built-in (non-LLM) table detector, exported to CSV/JSON
                so the "structured" parts of the site (course lists, faculty
                tables, fee schedules, merit lists, etc.) survive as clean
                tabular data instead of being flattened into prose.
     - PDFs / DOCX / XLSX / PPTX / CSV -> downloaded and parsed with
                dedicated non-LLM libraries (pdfplumber, python-docx,
                openpyxl, python-pptx, xlrd for legacy .xls).
2. Every artifact (page or file) gets a metadata record written next to it
   AND appended to a single JSONL manifest (`manifest.jsonl`) so a RAG
   ingestion job can just stream that file and know exactly what each chunk
   of text is, where it came from, and how to cite it.
3. Two-phase pipeline:
     Phase 1 — Deep-crawl every HTML page under *.nust.edu.pk (BFS),
               extracting clean markdown + tables + all discovered links.
               Pages blocked by anti-bot heuristics get one stealth retry
               (magic mode + user simulation) at the end of the phase.
     Phase 2 — Every discovered link to a document (pdf/doc/docx/xls/xlsx/
               ppt/pptx/csv/...) is downloaded once (dedup by URL + sha256)
               and parsed to text/tables. Downloads respect robots.txt,
               retry transient failures with backoff, enforce a size cap,
               and sanity-check content (an HTML login page served where a
               PDF was expected is recorded as a failure, not parsed).
4. RAG-oriented output:
     - Each page .md file starts with YAML frontmatter (source_url, title,
       fetched_at) so chunks stay citable even after the manifest is gone.
     - Markdown is post-cleaned: image tags, data: URIs, zero-width chars
       and blank-line runs are stripped.
     - Exact-duplicate page content (same body under two URLs) is detected
       via sha256 and recorded as `duplicate_of` so the ingester can skip it.
     - Scanned/image-only PDFs that yield no text are flagged `needs_ocr`.
     - Records with < 20 words are flagged `low_content`.
     - All paths in manifest/metadata are relative to the output dir, so the
       output folder can be moved/shipped as a unit.
5. Polite by default: respects robots.txt in BOTH phases, caps concurrency,
   adds a small delay between requests, identifies itself with a descriptive
   User-Agent, and only downloads documents hosted on *.nust.edu.pk unless
   --external-docs is given.
6. Resumable: both phases skip URLs that already have a completed record in
   the manifest (URL-canonicalized, so http/https/www variants count as the
   same page), so you can Ctrl+C and re-run safely.

Install
-------
    pip install -U crawl4ai httpx beautifulsoup4 lxml pandas \
                    pdfplumber python-docx python-pptx openpyxl aiofiles xlrd
    crawl4ai-setup        # one-time: installs the Playwright browser

Run
---
    python nust_crawler.py                      # crawl everything (careful!)
    python nust_crawler.py --max-pages 300 --max-depth 4    # a bounded test run
    python nust_crawler.py --start-url https://seecs.nust.edu.pk

Output layout
-------------
    output/
      manifest.jsonl              <- one JSON record per page/file (append-only)
      pages/<slug>.md             <- cleaned markdown per HTML page (with frontmatter)
      pages/<slug>.tables/*.csv   <- any <table> elements found on that page
      pages/<slug>.meta.json      <- per-page metadata
      files/<sha256_8>_<name>.pdf <- raw downloaded file (original bytes)
      files/<sha256_8>_<name>.txt <- extracted text
      files/<sha256_8>_<name>.tables/*.csv  <- tables extracted from the file
      files/<sha256_8>_<name>.meta.json     <- per-file metadata
      logs/crawl.log
"""

from __future__ import annotations

import argparse
import asyncio
import csv
import hashlib
import io
import json
import logging
import re
import sys
import time
import urllib.robotparser
from dataclasses import dataclass, field
from pathlib import Path
from urllib.parse import urlparse, urljoin, parse_qs
from datetime import datetime, timezone

# ---- crawl4ai ---------------------------------------------------------
from crawl4ai import AsyncWebCrawler, BrowserConfig, CrawlerRunConfig, CacheMode
from crawl4ai.deep_crawling import BFSDeepCrawlStrategy
from crawl4ai.deep_crawling.filters import (
    FilterChain,
    URLFilter,
    URLPatternFilter,
    ContentTypeFilter,
)
from crawl4ai.content_scraping_strategy import LXMLWebScrapingStrategy
from crawl4ai.content_filter_strategy import PruningContentFilter
from crawl4ai.markdown_generation_strategy import DefaultMarkdownGenerator
from crawl4ai.table_extraction import DefaultTableExtraction

# ---- third-party helpers -----------------------------------------------
import httpx

# =========================================================================
# Configuration
# =========================================================================

ROOT_DOMAIN = "nust.edu.pk"
DEFAULT_START_URL = "https://nust.edu.pk/"

# Extensions we treat as "documents" -> downloaded & parsed in Phase 2,
# and NOT followed as HTML pages during the deep crawl.
DOCUMENT_EXTENSIONS = {
    ".pdf", ".doc", ".docx", ".ppt", ".pptx", ".xls", ".xlsx",
    ".csv", ".rtf", ".odt", ".ods", ".odp", ".txt",
}
# Extensions we simply record (metadata only), but don't try to parse text from.
BINARY_EXTENSIONS = {
    ".zip", ".rar", ".7z", ".jpg", ".jpeg", ".png", ".gif", ".webp",
    ".mp4", ".mp3", ".avi", ".mov",
}

USER_AGENT = (
    "NUST-RAG-Crawler/1.1 (+research/indexing bot; "
    "contact: set-your-contact-email-here)"
)

# Records with fewer words than this get flagged `low_content` in the
# manifest so the RAG ingester can filter out empty shells / JS-only pages.
LOW_CONTENT_WORDS = 20

# =========================================================================
# Small utilities
# =========================================================================


def utcnow() -> str:
    return datetime.now(timezone.utc).isoformat()


def sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def sha256_text(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8", "ignore")).hexdigest()


def canonical_url(url: str) -> str:
    """Collapse trivially-equivalent URL variants into one key:
    scheme dropped, host lower-cased, leading www. stripped, trailing
    slash stripped. Used by the crawl-time dedupe filter AND the manifest
    resume check, so `http://nust.edu.pk/` never gets re-crawled after
    `https://nust.edu.pk/` succeeded.
    """
    p = urlparse(url)
    host = (p.netloc or "").lower()
    if host.startswith("www."):
        host = host[4:]
    path = p.path.rstrip("/") or "/"
    canon = f"{host}{path}"
    if p.query:
        canon += f"?{p.query}"
    return canon


def slugify_url(url: str, max_len: int = 150) -> str:
    """Turn a URL into a filesystem-safe, still-recognisable slug."""
    parsed = urlparse(url)
    base = f"{parsed.netloc}{parsed.path}"
    base = base.rstrip("/") or parsed.netloc
    base = re.sub(r"[^a-zA-Z0-9_.\-/]", "_", base)
    base = base.replace("/", "__")
    if parsed.query:
        base += "__q_" + sha256_text(parsed.query)[:10]
    if len(base) > max_len:
        base = base[:max_len] + "_" + sha256_text(url)[:10]
    return base or sha256_text(url)[:16]


def get_extension(url: str) -> str:
    path = urlparse(url).path.lower()
    for ext in sorted(DOCUMENT_EXTENSIONS | BINARY_EXTENSIONS, key=len, reverse=True):
        if path.endswith(ext):
            return ext
    return Path(path).suffix.lower()


def is_document_url(url: str) -> bool:
    return get_extension(url) in DOCUMENT_EXTENSIONS


def is_binary_url(url: str) -> bool:
    return get_extension(url) in BINARY_EXTENSIONS


def in_root_domain(url: str, root_domain: str = ROOT_DOMAIN) -> bool:
    host = (urlparse(url).netloc or "").lower().split(":")[0]
    return host == root_domain or host.endswith("." + root_domain)


def rel_path(path: Path, out_dir: Path) -> str:
    """Store paths relative to the output dir (posix-style) so the whole
    output folder can be moved between machines without breaking the
    manifest."""
    try:
        return path.resolve().relative_to(out_dir.resolve()).as_posix()
    except ValueError:
        return str(path)


# ---- markdown post-cleaning for RAG -------------------------------------

_MD_IMAGE_RE = re.compile(r"!\[[^\]]*\]\([^)]*\)")
_DATA_URI_RE = re.compile(r"\(data:[^)]{20,}\)")
_TRAILING_WS_RE = re.compile(r"[ \t]+\n")
_BLANK_RUNS_RE = re.compile(r"\n{3,}")


def clean_markdown(text: str) -> str:
    """Strip elements that add tokens but no retrievable meaning:
    image embeds, base64 data URIs, zero-width chars, trailing whitespace,
    runs of blank lines."""
    if not text:
        return ""
    text = _MD_IMAGE_RE.sub("", text)
    text = _DATA_URI_RE.sub("(image)", text)
    text = text.replace("​", "").replace("﻿", "").replace(" ", " ")
    text = _TRAILING_WS_RE.sub("\n", text)
    text = _BLANK_RUNS_RE.sub("\n\n", text)
    return text.strip()


def yaml_frontmatter(url: str, title: str, fetched_at: str) -> str:
    """Minimal YAML frontmatter so each .md file is self-describing/citable
    even when processed standalone by a chunker."""
    def esc(s: str) -> str:
        return (s or "").replace("\\", "\\\\").replace('"', '\\"')

    return (
        "---\n"
        f'source_url: "{esc(url)}"\n'
        f'title: "{esc(title)}"\n'
        f'fetched_at: "{fetched_at}"\n'
        "---\n\n"
    )


def short_error(msg: str, limit: int = 220) -> str:
    """crawl4ai sometimes returns multi-line error messages with an embedded
    source-code traceback (Playwright navigation errors etc). Keep the log
    readable; the full message is still stored in the per-record metadata.
    """
    if not msg:
        return ""
    first_line = msg.strip().splitlines()[0]
    return (first_line[:limit] + "…") if len(first_line) > limit else first_line


# =========================================================================
# Custom URL filters (subclassing crawl4ai's URLFilter, which just needs
# an `apply(self, url) -> bool` method)
# =========================================================================


class SubdomainFilter(URLFilter):
    """Allow the root domain and any of its subdomains; reject everything else.

    This matters for a university site like nust.edu.pk, where individual
    schools/departments (seecs.nust.edu.pk, sns.nust.edu.pk, s3h.nust.edu.pk,
    ...) are often served from separate subdomains that should still be
    considered part of "the entirety of the website".
    """

    def __init__(self, root_domain: str):
        super().__init__(name="SubdomainFilter")
        self.root_domain = root_domain.lower()

    def apply(self, url: str) -> bool:
        host = (urlparse(url).netloc or "").lower().split(":")[0]
        passed = host == self.root_domain or host.endswith("." + self.root_domain)
        self._update_stats(passed)
        return passed


class SkipDocumentExtensionsFilter(URLFilter):
    """Don't let the deep-crawl strategy try to render document files as HTML
    pages (that wastes a browser render + usually fails). We still capture
    these URLs separately via result.links for Phase 2 downloading.
    """

    def __init__(self):
        super().__init__(name="SkipDocumentExtensionsFilter")

    def apply(self, url: str) -> bool:
        passed = not (is_document_url(url) or is_binary_url(url))
        self._update_stats(passed)
        return passed


class DedupeFilter(URLFilter):
    """Reject URLs that are equivalent to one already accepted in this run.

    Real sites are full of near-duplicate links to the "same" page:
    http vs https, trailing slash vs no trailing slash, www vs no www.
    crawl4ai's own visited-set uses exact string matching, so
    `http://nust.edu.pk/` and `https://nust.edu.pk/` both get crawled as if
    they were different pages. This filter canonicalizes and only lets the
    first variant of each canonical URL through.

    The start URL never passes through the filter chain, so run_html_crawl
    seeds it via mark_seen() — otherwise the http:// variant of the start
    page gets crawled again at depth 2 (observed in real runs).
    """

    def __init__(self):
        super().__init__(name="DedupeFilter")
        self._seen: set[str] = set()

    def mark_seen(self, url: str):
        self._seen.add(canonical_url(url))

    def apply(self, url: str) -> bool:
        canon = canonical_url(url)
        passed = canon not in self._seen
        if passed:
            self._seen.add(canon)
        self._update_stats(passed)
        return passed


# Query-string keys that are near-guaranteed infinite-crawl traps: session
# IDs (a fresh value on every request -> every URL looks "new" forever),
# WordPress comment-reply permalinks (one per comment, recursively), and
# tracking params that don't change the page content at all.
TRAP_QUERY_KEYS = {
    "replytocom", "sessionid", "phpsessid", "jsessionid", "sid", "session",
    "utm_source", "utm_medium", "utm_campaign", "utm_term", "utm_content",
    "fbclid", "gclid", "share", "print",
}


class TrapQueryFilter(URLFilter):
    """Reject URLs carrying a known trap query parameter outright — this is
    a real, common cause of a crawler that never seems to finish: e.g.
    WordPress's `?replytocom=1234` appears on every comment of every post,
    each one technically a "new" URL, so a naive crawler can spend forever
    walking comment threads instead of actual site content.
    """

    def __init__(self):
        super().__init__(name="TrapQueryFilter")

    def apply(self, url: str) -> bool:
        query_keys = {k.lower() for k in parse_qs(urlparse(url).query).keys()}
        passed = not (query_keys & TRAP_QUERY_KEYS)
        self._update_stats(passed)
        return passed


class PathBudgetFilter(URLFilter):
    """General-purpose trap guard, independent of any specific parameter
    name: once a host+path has produced this many distinct query-string
    variants, stop accepting more of them.

    This is what actually protects against the "endlessly looping through
    generated links" symptom for cases TrapQueryFilter doesn't know about
    by name — calendar/event widgets (`?month=`, `?date=`), faceted search
    and filter UIs, or any other combinatorial-URL generator. A real page
    almost never needs 20+ query-string variants crawled to be represented
    in a RAG index; a trap can otherwise generate URLs indefinitely.
    """

    def __init__(self, max_variants_per_path: int = 20):
        super().__init__(name="PathBudgetFilter")
        self.max_variants = max_variants_per_path
        self._counts: dict[str, int] = {}

    def apply(self, url: str) -> bool:
        p = urlparse(url)
        if not p.query:
            self._update_stats(True)
            return True  # the bare path itself is never budget-limited
        key = f"{p.netloc.lower()}{p.path}"
        count = self._counts.get(key, 0)
        passed = count < self.max_variants
        if passed:
            self._counts[key] = count + 1
        self._update_stats(passed)
        return passed


class RepeatedPathSegmentFilter(URLFilter):
    """Reject URLs whose path repeats a segment, and cap path depth.

    Observed on nust.edu.pk: a bug in the faculty-profile page template
    renders publication links without a scheme — `href="doi.org/10.1007/x"`
    instead of `href="https://doi.org/10.1007/x"`. A scheme-less href is a
    *relative* link per the HTML spec, so it resolves onto the current
    profile URL's path instead of navigating to doi.org:

        .../faculty/name/doi.org/10.1007/A
        .../faculty/name/doi.org/10.1007/A/doi.org/10.1007/B
        .../faculty/name/doi.org/10.1007/A/doi.org/10.1007/B/doi.org/10.1007/C
        ...

    Because the malformed page repeats the same broken links, each
    additional publication nests deeper and the URL count grows
    combinatorially. A single 12-hour run of this crawler produced ~4,500
    such URLs across just 79 faculty profiles before being stopped.

    A legitimate path on this site never repeats a segment (no page is
    nested under itself), so "any repeated segment" is a clean, general
    signal for this entire bug class — it isn't specific to doi.org and
    would equally catch the same bug via linkedin.com, researchgate.net,
    orcid.org, or any other scheme-less external link on a template with
    this mistake. The max-depth cap is a belt-and-suspenders backstop for
    any other recursive-link bug that doesn't happen to repeat a segment.
    """

    def __init__(self, max_path_segments: int = 8):
        super().__init__(name="RepeatedPathSegmentFilter")
        self.max_path_segments = max_path_segments

    def apply(self, url: str) -> bool:
        segments = [s for s in urlparse(url).path.split("/") if s]
        passed = len(segments) <= self.max_path_segments and len(segments) == len(set(segments))
        self._update_stats(passed)
        return passed


# =========================================================================
# Manifest writer (append-only JSONL, safe under a single asyncio task)
# =========================================================================


class Manifest:
    """Append-only JSONL manifest + in-memory indexes used for resume and
    dedup:
      - _seen_urls: canonicalized URLs with a completed ("ok") record
      - _content_owner: page-markdown sha256 -> first URL that produced it
      - _file_sha_owner: downloaded-file sha256 -> first URL it came from
    """

    def __init__(self, path: Path):
        self.path = path
        self._lock = asyncio.Lock()
        self._seen_urls: set[str] = set()
        self._content_owner: dict[str, str] = {}
        self._file_sha_owner: dict[str, str] = {}
        if path.exists():
            with path.open("r", encoding="utf-8") as f:
                for line in f:
                    line = line.strip()
                    if not line:
                        continue
                    try:
                        rec = json.loads(line)
                    except json.JSONDecodeError:
                        continue
                    if rec.get("status") == "ok":
                        self._index(rec)

    def _index(self, rec: dict):
        self._seen_urls.add(canonical_url(rec["url"]))
        csha = rec.get("content_sha256")
        if csha and not rec.get("duplicate_of"):
            self._content_owner.setdefault(csha, rec["url"])
        fsha = rec.get("sha256")
        if fsha and not rec.get("duplicate_of"):
            self._file_sha_owner.setdefault(fsha, rec["url"])

    def already_done(self, url: str) -> bool:
        return canonical_url(url) in self._seen_urls

    def content_owner(self, content_sha: str) -> str | None:
        return self._content_owner.get(content_sha)

    def file_sha_owner(self, file_sha: str) -> str | None:
        return self._file_sha_owner.get(file_sha)

    async def write(self, record: dict):
        async with self._lock:
            with self.path.open("a", encoding="utf-8") as f:
                f.write(json.dumps(record, ensure_ascii=False) + "\n")
            if record.get("status") == "ok":
                self._index(record)


# =========================================================================
# Robots.txt for Phase 2 (Phase 1 is handled by crawl4ai itself)
# =========================================================================


class RobotsCache:
    """Tiny per-host robots.txt cache for the document-download phase.
    Fails open: if robots.txt can't be fetched, downloads are allowed."""

    def __init__(self, user_agent: str):
        self._parsers: dict[str, urllib.robotparser.RobotFileParser] = {}
        self._lock = asyncio.Lock()
        self.user_agent = user_agent

    async def allowed(self, client: httpx.AsyncClient, url: str) -> bool:
        p = urlparse(url)
        origin = f"{p.scheme}://{p.netloc}"
        async with self._lock:
            rp = self._parsers.get(origin)
            if rp is None:
                rp = urllib.robotparser.RobotFileParser()
                try:
                    resp = await client.get(
                        origin + "/robots.txt", timeout=15.0, follow_redirects=True
                    )
                    if resp.status_code == 200:
                        rp.parse(resp.text.splitlines())
                    else:
                        rp.parse([])
                except Exception:
                    rp.parse([])
                self._parsers[origin] = rp
        return rp.can_fetch(self.user_agent, url)


# =========================================================================
# Table export helper (used for both HTML pages and parsed documents)
# =========================================================================


def is_meaningful_table(headers, rows, min_cols: int = 2, min_cells: int = 4) -> bool:
    """Filter out junk table fragments.

    pdfplumber (and, less often, crawl4ai's HTML table detector) regularly
    mis-detects a single rotated/merged label cell as its own "table" on
    borderless, calendar-grid-style layout PDFs — e.g. NUST's academic
    schedule PDF yields real tables plus 3 spurious one-cell "tables" that
    are each just a stray "1 Green week" label pulled out of a merged
    calendar cell. These carry no real tabular information (the same text
    is already present in the full extracted .txt/.md), so dropping them
    costs nothing and keeps the RAG index from picking up single-cell
    noise chunks and cuts down on clutter in files/*.tables/.

    Real tables (a fee schedule, a course list, ...) comfortably clear a
    "at least 2 columns and at least 4 total cells" bar.
    """
    n_cols = len(headers) if headers else (len(rows[0]) if rows else 0)
    n_rows_total = (1 if headers else 0) + len(rows)
    n_cells = n_cols * n_rows_total
    return n_cols >= min_cols and n_cells >= min_cells


def dump_tables(tables: list, out_dir: Path, stem: str, root_dir: Path) -> list[dict]:
    """tables: list of {"headers": [...], "rows": [[...], ...]} (crawl4ai's
    native table format) OR plain list-of-lists. Returns metadata list with
    paths relative to root_dir. Junk single-cell/single-column fragments
    (see is_meaningful_table) are dropped rather than written to disk.
    """
    if not tables:
        return []
    out_dir.mkdir(parents=True, exist_ok=True)
    written = []
    kept_index = 0
    for table in tables:
        headers = table.get("headers") if isinstance(table, dict) else None
        rows = table.get("rows") if isinstance(table, dict) else table
        if not rows:
            continue
        if not is_meaningful_table(headers, rows):
            continue
        i = kept_index
        kept_index += 1
        csv_path = out_dir / f"{stem}_table_{i}.csv"
        with csv_path.open("w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            if headers:
                writer.writerow(headers)
            for row in rows:
                writer.writerow(row)
        meta = {
            "table_index": i,
            "path": rel_path(csv_path, root_dir),
            "n_rows": len(rows),
            "n_cols": len(headers) if headers else (len(rows[0]) if rows else 0),
            "has_headers": bool(headers),
        }
        # carry through source context (page number / sheet / slide) if the
        # parser provided it — lets the RAG layer cite "page 3 of the PDF".
        for key in ("page", "sheet", "slide"):
            if isinstance(table, dict) and table.get(key) is not None:
                meta[key] = table[key]
        written.append(meta)
    return written


# =========================================================================
# Phase 1: HTML deep crawl
# =========================================================================


@dataclass
class CrawlStats:
    pages_ok: int = 0
    pages_failed: int = 0
    pages_duplicate: int = 0
    document_links_found: set = field(default_factory=set)


def make_markdown_generator() -> DefaultMarkdownGenerator:
    return DefaultMarkdownGenerator(
        # Non-LLM heuristic boilerplate remover (nav/footer/ads stripping).
        content_filter=PruningContentFilter(threshold=0.45, threshold_type="fixed"),
        options={"ignore_links": False, "citations": False},
    )


def build_run_config(
    max_depth: int,
    max_pages: int,
    dedupe: DedupeFilter,
    trap_query: TrapQueryFilter,
    path_budget: PathBudgetFilter,
    repeated_segment: RepeatedPathSegmentFilter,
) -> CrawlerRunConfig:
    filter_chain = FilterChain(
        [
            SubdomainFilter(ROOT_DOMAIN),
            SkipDocumentExtensionsFilter(),
            trap_query,        # cheap, explicit: known trap params (replytocom, sid, ...)
            repeated_segment,  # kills the doi.org-style relative-link recursion bug
            path_budget,       # general trap guard: caps query-variant explosion per path
            ContentTypeFilter(allowed_types=["text/html"]),
            dedupe,  # must come last: only records URLs that already
            # passed every other check, so a rejected doc-link doesn't
            # "use up" its canonical slot.
        ]
    )
    deep_crawl_strategy = BFSDeepCrawlStrategy(

        max_depth=max_depth,
        max_pages=max_pages,
        filter_chain=filter_chain,
        include_external=True,  # subdomains are technically a different host;
        # the SubdomainFilter above is what actually restricts scope.
    )
    return CrawlerRunConfig(
        deep_crawl_strategy=deep_crawl_strategy,
        scraping_strategy=LXMLWebScrapingStrategy(),
        markdown_generator=make_markdown_generator(),
        # Lowered from the library default (7) — the default threshold drops
        # small-but-real data tables (a 2-column fee table, a short contact
        # list) that matter for a university site with lots of small tables.
        table_extraction=DefaultTableExtraction(table_score_threshold=5),
        cache_mode=CacheMode.BYPASS,
        check_robots_txt=True,
        user_agent=USER_AGENT,
        exclude_external_links=False,  # we WANT to see doc links even if flagged external
        # Rewrite internal http:// links to https:// before they're queued,
        # so "http://nipo.nust.edu.pk/" and "https://nipo.nust.edu.pk/"
        # collapse into a single crawl instead of two.
        preserve_https_for_internal_links=True,
        stream=True,
        verbose=True,
        mean_delay=0.5,   # politeness delay between requests (seconds)
        max_range=1.0,
        semaphore_count=4,  # cap concurrent browser tabs
        page_timeout=45000,
        # A little extra settle time before we snapshot the HTML — helps
        # avoid false "anti-bot / empty shell" detections on pages that are
        # legitimately JS-rendered but just need a beat longer to mount.
        delay_before_return_html=1.0,
    )


def build_stealth_retry_config() -> CrawlerRunConfig:
    """Single-page config for retrying pages that failed the first pass with
    an anti-bot / empty-shell detection. magic + simulate_user are crawl4ai's
    non-LLM stealth heuristics; the longer settle time helps JS-heavy shells
    (observed: alumni.nust.edu.pk rendered nothing in 1s)."""
    return CrawlerRunConfig(
        scraping_strategy=LXMLWebScrapingStrategy(),
        markdown_generator=make_markdown_generator(),
        table_extraction=DefaultTableExtraction(table_score_threshold=5),
        cache_mode=CacheMode.BYPASS,
        check_robots_txt=True,
        user_agent=USER_AGENT,
        exclude_external_links=False,
        magic=True,
        simulate_user=True,
        override_navigator=True,
        page_timeout=60000,
        delay_before_return_html=4.0,
    )


async def process_page_result(
    result,
    out_dir: Path,
    pages_dir: Path,
    manifest: Manifest,
    stats: CrawlStats,
    logger: logging.Logger,
    retried: bool = False,
) -> None:
    """Write markdown + tables + metadata for one successful crawl result.
    Shared by the main BFS stream and the stealth retry pass."""
    url = result.url
    depth = (result.metadata or {}).get("depth", 0)
    fetched_at = utcnow()

    # ---- collect document links for phase 2 --------------------
    all_links = []
    for bucket in ("internal", "external"):
        all_links.extend(result.links.get(bucket, []) or [])
    doc_links_here = set()
    for link in all_links:
        href = link.get("href") if isinstance(link, dict) else link
        if not href:
            continue
        href = urljoin(url, href)
        if is_document_url(href) or is_binary_url(href):
            doc_links_here.add(href)
    stats.document_links_found |= doc_links_here

    # ---- clean markdown ----------------------------------------
    md_obj = result.markdown
    # fit_markdown = post PruningContentFilter (cleaner); fall back to raw.
    markdown_text = (
        getattr(md_obj, "fit_markdown", None)
        or getattr(md_obj, "raw_markdown", None)
        or str(md_obj or "")
    )
    markdown_text = clean_markdown(markdown_text)
    content_sha = sha256_text(markdown_text)
    word_count = len(markdown_text.split())

    # Exact-content dedup: same body already stored under another URL
    # (e.g. /counselling-services vs /counselling-services/, or an alias
    # subdomain). We still record the URL for citation purposes but point
    # at the owner instead of writing a second copy.
    duplicate_of = None
    if word_count > 0:
        owner = manifest.content_owner(content_sha)
        if owner and canonical_url(owner) != canonical_url(url):
            duplicate_of = owner

    title = (result.metadata or {}).get("title") or ""
    slug = slugify_url(url)
    md_path = pages_dir / f"{slug}.md"
    tables_meta: list[dict] = []
    if duplicate_of is None:
        md_path.write_text(
            yaml_frontmatter(url, title, fetched_at) + markdown_text + "\n",
            encoding="utf-8",
        )
        tables_meta = dump_tables(
            result.tables or [], pages_dir / f"{slug}.tables", slug, out_dir
        )

    record = {
        "type": "page",
        "status": "ok",
        "url": url,
        "redirected_url": result.redirected_url or url,
        "status_code": result.status_code,
        "depth": depth,
        "title": title,
        "fetched_at": fetched_at,
        "content_path": rel_path(md_path, out_dir) if duplicate_of is None else None,
        "content_sha256": content_sha,
        "word_count": word_count,
        "low_content": word_count < LOW_CONTENT_WORDS,
        "duplicate_of": duplicate_of,
        "retried_stealth": retried,
        "n_tables": len(tables_meta),
        "tables": tables_meta,
        "n_internal_links": len(result.links.get("internal", []) or []),
        "n_external_links": len(result.links.get("external", []) or []),
        "n_document_links_found_here": len(doc_links_here),
        "document_links_found_here": sorted(doc_links_here),
        "n_images": len(result.media.get("images", []) or []),
    }
    if duplicate_of is None:
        meta_path = pages_dir / f"{slug}.meta.json"
        meta_path.write_text(
            json.dumps(record, indent=2, ensure_ascii=False), encoding="utf-8"
        )

    await manifest.write(record)
    if duplicate_of:
        stats.pages_duplicate += 1
        logger.info(f"[ok depth={depth}, duplicate of {duplicate_of}] {url}")
    else:
        stats.pages_ok += 1
        flag = " LOW-CONTENT" if record["low_content"] else ""
        logger.info(
            f"[ok depth={depth}{flag}] {url}  "
            f"({word_count} words, {len(tables_meta)} tables)"
        )


async def run_html_crawl(
    start_url: str,
    out_dir: Path,
    manifest: Manifest,
    max_depth: int,
    max_pages: int,
    logger: logging.Logger,
    retry_blocked: bool = True,
) -> CrawlStats:
    stats = CrawlStats()
    pages_dir = out_dir / "pages"
    pages_dir.mkdir(parents=True, exist_ok=True)

    browser_config = BrowserConfig(headless=True, user_agent=USER_AGENT, verbose=False)
    dedupe = DedupeFilter()
    dedupe.mark_seen(start_url)  # start URL bypasses the filter chain
    trap_query = TrapQueryFilter()
    path_budget = PathBudgetFilter()
    repeated_segment = RepeatedPathSegmentFilter()
    run_config = build_run_config(
        max_depth, max_pages, dedupe, trap_query, path_budget, repeated_segment
    )

    failed: dict[str, str] = {}  # url -> error (candidates for stealth retry)

    async with AsyncWebCrawler(config=browser_config) as crawler:
        stream = await crawler.arun(start_url, config=run_config)
        async for result in stream:
            url = result.url
            depth = (result.metadata or {}).get("depth", 0)

            if manifest.already_done(url):
                logger.info(f"[skip: already done] {url}")
                continue

            if not result.success:
                logger.warning(f"[FAILED] {url}: {short_error(result.error_message)}")
                failed[url] = result.error_message or ""
                continue

            await process_page_result(result, out_dir, pages_dir, manifest, stats, logger)

        # ---- diagnostics: did we hit any crawl-trap guards? -------------
        # Non-zero rejected counts here are the concrete evidence for
        # whether the crawl was fighting a trap (session IDs, comment-reply
        # permalinks, calendar/faceted-search query explosion) rather than
        # genuinely working through new site content.
        for f in (dedupe, trap_query, path_budget, repeated_segment):
            if f.stats.rejected_urls:
                logger.info(
                    f"[trap-guard] {f.name}: rejected {f.stats.rejected_urls} / "
                    f"{f.stats.total_urls} candidate URLs"
                )

        # ---- stealth retry pass for blocked / empty-shell pages --------
        if retry_blocked and failed:
            logger.info(f"Retrying {len(failed)} failed page(s) with stealth config...")
            stealth_config = build_stealth_retry_config()
            for url, first_error in failed.items():
                try:
                    result = await crawler.arun(url, config=stealth_config)
                except Exception as e:
                    result = None
                    retry_error = str(e)
                if result is not None and result.success:
                    await process_page_result(
                        result, out_dir, pages_dir, manifest, stats, logger, retried=True
                    )
                    continue
                retry_error = (
                    result.error_message if result is not None else retry_error
                )
                stats.pages_failed += 1
                logger.warning(f"[FAILED after retry] {url}: {short_error(retry_error)}")
                await manifest.write(
                    {
                        "type": "page",
                        "url": url,
                        "status": "failed",
                        "error": first_error,
                        "retry_error": retry_error,
                        "fetched_at": utcnow(),
                    }
                )
        else:
            for url, err in failed.items():
                stats.pages_failed += 1
                await manifest.write(
                    {
                        "type": "page",
                        "url": url,
                        "status": "failed",
                        "error": err,
                        "fetched_at": utcnow(),
                    }
                )

    return stats


# =========================================================================
# Phase 2: document download + parse
# =========================================================================


def parse_pdf(data: bytes) -> tuple[str, list, dict]:
    import pdfplumber

    text_parts = []
    tables = []
    extra = {}
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        extra["n_pages"] = len(pdf.pages)
        for i, page in enumerate(pdf.pages):
            page_text = page.extract_text() or ""
            if page_text:
                text_parts.append(f"\n\n--- Page {i + 1} ---\n{page_text}")
            for tbl in page.extract_tables() or []:
                if tbl:
                    tables.append({"headers": tbl[0], "rows": tbl[1:], "page": i + 1})
    return "".join(text_parts).strip(), tables, extra


def parse_docx(data: bytes) -> tuple[str, list, dict]:
    import docx

    doc = docx.Document(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    tables = []
    for t in doc.tables:
        rows = [[cell.text for cell in row.cells] for row in t.rows]
        if rows:
            tables.append({"headers": rows[0], "rows": rows[1:]})
    text = "\n\n".join(paragraphs)
    return text, tables, {"n_paragraphs": len(paragraphs), "n_tables": len(doc.tables)}


def parse_xlsx(data: bytes) -> tuple[str, list, dict]:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    tables = []
    text_parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = [[("" if c is None else str(c)) for c in row] for row in ws.iter_rows(values_only=True)]
        rows = [r for r in rows if any(cell.strip() for cell in r)]
        if not rows:
            continue
        tables.append({"headers": rows[0], "rows": rows[1:], "sheet": sheet_name})
        text_parts.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(", ".join(r) for r in rows))
    return "\n\n".join(text_parts), tables, {"n_sheets": len(wb.sheetnames), "sheet_names": wb.sheetnames}


def parse_xls(data: bytes) -> tuple[str, list, dict]:
    # openpyxl can NOT read legacy binary .xls — xlrd can. (Routing .xls to
    # openpyxl was a silent parse-failure in the previous version.)
    import xlrd

    book = xlrd.open_workbook(file_contents=data)
    tables = []
    text_parts = []
    for sheet in book.sheets():
        rows = [
            [("" if c is None else str(c)) for c in sheet.row_values(r)]
            for r in range(sheet.nrows)
        ]
        rows = [r for r in rows if any(cell.strip() for cell in r)]
        if not rows:
            continue
        tables.append({"headers": rows[0], "rows": rows[1:], "sheet": sheet.name})
        text_parts.append(f"--- Sheet: {sheet.name} ---\n" + "\n".join(", ".join(r) for r in rows))
    return "\n\n".join(text_parts), tables, {"n_sheets": book.nsheets}


def parse_pptx(data: bytes) -> tuple[str, list, dict]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    text_parts = []
    tables = []
    n_slides = 0
    for i, slide in enumerate(prs.slides):
        n_slides += 1
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        slide_text.append(line)
            if shape.has_table:
                rows = [
                    [cell.text for cell in row.cells] for row in shape.table.rows
                ]
                if rows:
                    tables.append({"headers": rows[0], "rows": rows[1:], "slide": i + 1})
        if slide_text:
            text_parts.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_text))
    return "\n\n".join(text_parts), tables, {"n_slides": n_slides}


def parse_csv_bytes(data: bytes) -> tuple[str, list, dict]:
    text = data.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = list(reader)
    tables = []
    if rows:
        tables.append({"headers": rows[0], "rows": rows[1:]})
    return text, tables, {"n_rows": len(rows)}


PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".xlsx": parse_xlsx,
    ".xls": parse_xls,
    ".pptx": parse_pptx,
    ".csv": parse_csv_bytes,
    ".txt": lambda data: (data.decode("utf-8", "replace"), [], {}),
}

# Magic-byte signatures used to catch servers that answer a document URL
# with an HTML error/login page (would otherwise fail the parser with a
# confusing error, or worse, get indexed as garbage text).
_MAGIC_SIGNATURES = {
    ".pdf": (b"%PDF",),
    ".docx": (b"PK\x03\x04",),
    ".xlsx": (b"PK\x03\x04",),
    ".pptx": (b"PK\x03\x04",),
    ".doc": (b"\xd0\xcf\x11\xe0",),
    ".xls": (b"\xd0\xcf\x11\xe0",),
    ".ppt": (b"\xd0\xcf\x11\xe0",),
}


def content_mismatch(ext: str, data: bytes, content_type: str | None) -> str | None:
    """Return a reason string if the downloaded bytes clearly aren't the
    file type the URL promised; None if they look fine."""
    sigs = _MAGIC_SIGNATURES.get(ext)
    if sigs and not any(data.startswith(s) for s in sigs):
        head = data[:256].lstrip().lower()
        if head.startswith(b"<!doctype") or head.startswith(b"<html"):
            return f"expected {ext} but server returned an HTML page"
        return f"expected {ext} but magic bytes do not match"
    return None


async def fetch_with_retries(
    client: httpx.AsyncClient,
    url: str,
    max_bytes: int,
    attempts: int = 3,
) -> httpx.Response | tuple[None, str]:
    """Streamed GET with a size cap and simple exponential backoff on
    transient errors (network hiccups, 5xx, 429). Returns the response with
    .content populated, or (None, error_string)."""
    last_error = ""
    for attempt in range(attempts):
        if attempt:
            await asyncio.sleep(2 ** attempt)
        try:
            async with client.stream(
                "GET", url, timeout=60.0, follow_redirects=True
            ) as resp:
                if resp.status_code in (429, 500, 502, 503, 504):
                    last_error = f"HTTP {resp.status_code}"
                    continue
                resp.raise_for_status()
                clen = resp.headers.get("content-length")
                if clen and int(clen) > max_bytes:
                    return None, f"file too large ({clen} bytes > cap {max_bytes})"
                chunks = []
                total = 0
                async for chunk in resp.aiter_bytes():
                    total += len(chunk)
                    if total > max_bytes:
                        return None, f"file too large (> cap {max_bytes} bytes)"
                    chunks.append(chunk)
                resp._content = b"".join(chunks)
                return resp
        except httpx.HTTPStatusError as e:
            # 4xx (other than 429) won't get better on retry.
            return None, str(e)
        except httpx.HTTPError as e:
            last_error = str(e)
    return None, f"failed after {attempts} attempts: {last_error}"


async def download_and_parse(
    url: str,
    referring_pages: list[str],
    client: httpx.AsyncClient,
    out_dir: Path,
    manifest: Manifest,
    logger: logging.Logger,
    sem: asyncio.Semaphore,
    robots: RobotsCache,
    max_bytes: int,
):
    if manifest.already_done(url):
        logger.info(f"[skip file, already done] {url}")
        return

    files_dir = out_dir / "files"
    files_dir.mkdir(parents=True, exist_ok=True)

    ext = get_extension(url)
    record = {
        "type": "file",
        "status": "failed",
        "url": url,
        "extension": ext,
        "referring_pages": referring_pages,
        "fetched_at": utcnow(),
    }

    async with sem:
        if not await robots.allowed(client, url):
            record["error"] = "disallowed by robots.txt"
            await manifest.write(record)
            logger.info(f"[skip file, robots.txt] {url}")
            return
        fetched = await fetch_with_retries(client, url, max_bytes)
        if isinstance(fetched, tuple):
            record["error"] = fetched[1]
            await manifest.write(record)
            logger.warning(f"[FAILED download] {url}: {short_error(fetched[1])}")
            return
        resp = fetched
        data = resp.content

    digest = sha256_bytes(data)
    record.update(
        {
            "content_type_header": resp.headers.get("content-type"),
            "size_bytes": len(data),
            "sha256": digest,
            "final_url": str(resp.url),
        }
    )

    # Same bytes already downloaded under a different URL (mirrored file):
    # record the alias, don't store or parse a second copy.
    owner = manifest.file_sha_owner(digest)
    if owner and canonical_url(owner) != canonical_url(url):
        record.update({"status": "ok", "parsed": False, "duplicate_of": owner})
        await manifest.write(record)
        logger.info(f"[ok, duplicate of {owner}] {url}")
        return

    # Catch HTML error/login pages served where a document was promised.
    mismatch = content_mismatch(ext, data, resp.headers.get("content-type"))
    if mismatch:
        record["error"] = mismatch
        await manifest.write(record)
        logger.warning(f"[FAILED, content mismatch] {url}: {mismatch}")
        return

    original_name = Path(urlparse(url).path).name or digest[:16]
    safe_name = re.sub(r"[^a-zA-Z0-9_.\-]", "_", original_name)
    base_name = f"{digest[:10]}_{safe_name}"
    raw_path = files_dir / base_name
    raw_path.write_bytes(data)
    record["raw_path"] = rel_path(raw_path, out_dir)

    parser = PARSERS.get(ext)
    if parser is None:
        # binary / unparsed file: keep metadata only
        record["status"] = "ok"
        record["parsed"] = False
        await manifest.write(record)
        logger.info(f"[ok, unparsed binary] {url} ({len(data)} bytes)")
        return

    try:
        text, tables, extra = parser(data)
    except Exception as e:
        record["status"] = "ok"
        record["parsed"] = False
        record["parse_error"] = str(e)
        await manifest.write(record)
        logger.warning(f"[downloaded but PARSE FAILED] {url}: {e}")
        return

    text = (text or "").strip()
    word_count = len(text.split())
    text_path = files_dir / f"{base_name}.txt"
    text_path.write_text(text, encoding="utf-8")
    tables_meta = dump_tables(tables, files_dir / f"{base_name}.tables", base_name, out_dir)

    # A PDF that parses "successfully" to zero text is almost always a
    # scanned/image-only document — flag it so the RAG pipeline can route
    # it to OCR instead of silently indexing an empty string.
    needs_ocr = ext == ".pdf" and word_count == 0 and not tables_meta

    record.update(
        {
            "status": "ok",
            "parsed": True,
            "text_path": rel_path(text_path, out_dir),
            "content_sha256": sha256_text(text),
            "word_count": word_count,
            "low_content": word_count < LOW_CONTENT_WORDS,
            "needs_ocr": needs_ocr,
            "n_tables": len(tables_meta),
            "tables": tables_meta,
            "extra": extra,
        }
    )
    meta_path = files_dir / f"{base_name}.meta.json"
    meta_path.write_text(json.dumps(record, indent=2, ensure_ascii=False), encoding="utf-8")
    await manifest.write(record)
    flag = " NEEDS-OCR" if needs_ocr else ""
    logger.info(
        f"[ok, parsed {ext}{flag}] {url} ({word_count} words, {len(tables_meta)} tables)"
    )


async def run_document_phase(
    document_links: dict[str, list[str]],  # url -> [referring pages]
    out_dir: Path,
    manifest: Manifest,
    logger: logging.Logger,
    concurrency: int = 12,
    max_file_mb: int = 50,
):
    sem = asyncio.Semaphore(concurrency)
    robots = RobotsCache(USER_AGENT)
    max_bytes = max_file_mb * 1024 * 1024
    headers = {"User-Agent": USER_AGENT}
    # NOTE: previously this created each coroutine, scheduled it once via
    # asyncio.ensure_future(), then handed the *same* coroutine objects to
    # asyncio.gather(), which schedules them a second time. That double
    # scheduling caused httpx's internal anyio cancel-scopes to be entered
    # in one Task and exited in another -> "Attempted to exit cancel scope
    # in a different task than it was entered in", and every download in
    # the batch failed. Fix: create each Task exactly once (via
    # asyncio.create_task), stagger creation for politeness, then await
    # all of them together.
    async with httpx.AsyncClient(headers=headers) as client:
        tasks = []
        for i, (url, refs) in enumerate(document_links.items()):
            if i and i % concurrency == 0:
                await asyncio.sleep(1.0)
            tasks.append(
                asyncio.create_task(
                    download_and_parse(
                        url, refs, client, out_dir, manifest, logger, sem,
                        robots, max_bytes,
                    )
                )
            )
        if tasks:
            await asyncio.gather(*tasks)


# =========================================================================
# Main
# =========================================================================


def build_logger(out_dir: Path) -> logging.Logger:
    logs_dir = out_dir / "logs"
    logs_dir.mkdir(parents=True, exist_ok=True)
    logger = logging.getLogger("nust_crawler")
    logger.setLevel(logging.INFO)
    logger.handlers.clear()

    fh = logging.FileHandler(logs_dir / "crawl.log", encoding="utf-8")
    fh.setFormatter(logging.Formatter("%(asctime)s %(levelname)s %(message)s"))
    logger.addHandler(fh)

    sh = logging.StreamHandler(sys.stdout)
    sh.setFormatter(logging.Formatter("%(message)s"))
    logger.addHandler(sh)
    return logger


async def main_async(args):
    out_dir = Path(args.output)
    out_dir.mkdir(parents=True, exist_ok=True)
    logger = build_logger(out_dir)
    manifest = Manifest(out_dir / "manifest.jsonl")

    stats = None
    if args.skip_html:
        pages_dir = out_dir / "pages"
        if not pages_dir.exists() or not any(pages_dir.glob("*.meta.json")):
            raise SystemExit(
                f"--skip-html given but no page metadata found under {pages_dir}. "
                "Phase 1 needs to have completed at least once before you can jump "
                "straight to Phase 2."
            )
        logger.info("=" * 70)
        logger.info("Skipping Phase 1 (--skip-html given) — resuming straight into "
                     "Phase 2 using the page metadata already on disk. No pages will "
                     "be re-fetched or re-hashed.")
        logger.info("=" * 70)
    else:
        logger.info("=" * 70)
        logger.info(f"PHASE 1: HTML deep crawl starting at {args.start_url}")
        logger.info(f"max_depth={args.max_depth}  max_pages={args.max_pages}")
        logger.info("=" * 70)

        t0 = time.time()
        stats = await run_html_crawl(
            args.start_url, out_dir, manifest, args.max_depth, args.max_pages, logger,
            retry_blocked=not args.no_retry_blocked,
        )
        logger.info(
            f"Phase 1 done in {time.time() - t0:.1f}s: "
            f"{stats.pages_ok} pages ok, {stats.pages_failed} failed, "
            f"{stats.pages_duplicate} duplicates, "
            f"{len(stats.document_links_found)} document links discovered."
        )

    if args.skip_files:
        logger.info("Skipping Phase 2 (--skip-files given).")
        return

    # Build url -> referring pages map by re-reading page metadata just
    # written. Reading meta files (not just this run's stats) means a
    # resumed run — or a --skip-html run — still downloads documents
    # discovered by earlier runs, without needing Phase 1's in-memory stats.
    referring: dict[str, list[str]] = {}
    for meta_file in (out_dir / "pages").glob("*.meta.json"):
        try:
            rec = json.loads(meta_file.read_text(encoding="utf-8"))
        except json.JSONDecodeError:
            continue
        for link in rec.get("document_links_found_here", []):
            referring.setdefault(link, []).append(rec["url"])

    # Make sure every discovered doc link is present even if the referring
    # map build above missed something (defensive). Only relevant when
    # Phase 1 actually ran this invocation — with --skip-html, the meta
    # files above are the sole (and authoritative) source.
    if stats is not None:
        for link in stats.document_links_found:
            referring.setdefault(link, [])

    # By default only download documents hosted on *.nust.edu.pk — pages
    # link out to third-party PDFs (gov portals etc.) that aren't "the
    # website" and often 404. Opt in with --external-docs.
    if not args.external_docs:
        external = [u for u in referring if not in_root_domain(u)]
        for u in external:
            referring.pop(u)
        if external:
            logger.info(
                f"Skipping {len(external)} off-domain document link(s) "
                "(use --external-docs to include them)."
            )

    logger.info("=" * 70)
    logger.info(f"PHASE 2: downloading & parsing {len(referring)} document(s)")
    logger.info("=" * 70)

    t1 = time.time()
    await run_document_phase(
        referring, out_dir, manifest, logger,
        concurrency=args.file_concurrency,
        max_file_mb=args.max_file_mb,
    )
    logger.info(f"Phase 2 done in {time.time() - t1:.1f}s.")

    logger.info("=" * 70)
    logger.info(f"ALL DONE. Manifest: {out_dir / 'manifest.jsonl'}")
    logger.info("=" * 70)


def parse_args():
    p = argparse.ArgumentParser(description="Crawl & scrape nust.edu.pk for a RAG pipeline.")
    p.add_argument("--start-url", default=DEFAULT_START_URL)
    p.add_argument("--output", default="output")
    p.add_argument("--max-depth", type=int, default=6, help="BFS depth from the start URL.")
    p.add_argument(
        "--max-pages",
        type=int,
        default=100000,
        help="Safety cap on number of HTML pages to crawl. Set lower for a test run.",
    )
    p.add_argument(
        "--skip-files",
        action="store_true",
        help="Only do the HTML crawl; skip downloading/parsing PDFs etc.",
    )
    p.add_argument(
        "--skip-html",
        action="store_true",
        help=(
            "Skip Phase 1 entirely and go straight to Phase 2, using the "
            "page metadata already on disk from a previous run. Use this "
            "when Phase 1 already finished (or you have enough of it) and "
            "you just need to (re)do document downloading/parsing — no "
            "pages get re-fetched, re-scraped, or re-hashed."
        ),
    )
    p.add_argument(
        "--file-concurrency",
        type=int,
        default=5,
        help="Concurrent document downloads in Phase 2.",
    )
    p.add_argument(
        "--max-file-mb",
        type=int,
        default=50,
        help="Skip documents larger than this many MB in Phase 2.",
    )
    p.add_argument(
        "--external-docs",
        action="store_true",
        help="Also download documents hosted outside *.nust.edu.pk.",
    )
    p.add_argument(
        "--no-retry-blocked",
        action="store_true",
        help="Skip the stealth retry pass for pages blocked by anti-bot detection.",
    )
    return p.parse_args()


if __name__ == "__main__":
    asyncio.run(main_async(parse_args()))
